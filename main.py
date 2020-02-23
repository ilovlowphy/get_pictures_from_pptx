from pptx import Presentation, shapes
import os, re

def get_img_name_from_string(text):
    return re.sub('[\/:*?"<>|]','-',text) #去掉非法字符

def get_pictures_from_pptx(pptx):
    current_dir, pptx_name = os.path.split(pptx)
    img_dir = os.path.join(current_dir,os.path.splitext(pptx_name)[0])
    if not os.path.exists(img_dir):
        os.makedirs(img_dir)
    prs = Presentation(pptx)
    for i, slide in enumerate(prs.slides):
        page = i+1
        texts = []
        imgs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        texts.append(run.text)
            elif isinstance(shape, shapes.picture.Picture):
                imgs.append(shape.image)
        if not imgs:
            continue
        if texts:
            filename = get_img_name_from_string('_'.join(texts))[:20]
        else:
            filename = 'slide_'+str(page)
        if len(imgs) ==1:
            full_file_name = os.path.join(img_dir,filename+ '.'+ imgs[0].ext)
            if os.path.exists(full_file_name):
                full_file_name = os.path.join(img_dir,filename+ '_slide_'+str(page)+'.'+ imgs[0].ext)
            with open(full_file_name,'wb') as f:
                f.write(imgs[0].blob)
        else:
            for i, img in enumerate(imgs):
                full_file_name = os.path.join(img_dir,filename+'_'+ str(i)+ '.'+ imgs[0].ext)
                if os.path.exists(full_file_name):
                    full_file_name = os.path.join(img_dir,filename+'_'+ str(i)+'_slide_'+str(page)+ '.'+ imgs[0].ext)
                with open(full_file_name,'wb') as f:
                    f.write(img.blob)


if __name__ == '__main__':
    current_dir = os.sys.path[0]
    print(current_dir)
    dir_list = os.listdir(current_dir)
    ppt_file_names = (fns for fns in dir_list if fns.endswith('.pptx'))
    ppt_names = (os.path.splitext(fns)[0] for fns in dir_list if fns.endswith('.pptx'))
    for pptx in ppt_file_names:
        get_pictures_from_pptx(pptx)
